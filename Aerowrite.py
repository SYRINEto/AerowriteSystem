"""
AeroWrite — Tableau Blanc Interactif
- Mode TABLEAU BLANC  : fond blanc pur (ou slide PDF/PPTX), visage NON visible
- Mode CAMÉRA         : flux caméra en direct avec squelette main 21 points
- Suivi main en temps réel : 21 points de repère + connexions
- Dessin par gestes (MediaPipe Hands)
- Ouverture PDF / PowerPoint pour annoter
- Enregistrement vidéo + audio
Compatible MediaPipe 0.10+
"""

import cv2
import numpy as np
import time
import os
import threading
import urllib.request
import mediapipe as mp
from mediapipe.tasks import python as mp_python
from mediapipe.tasks.python import vision as mp_vision

APP_NAME = "AeroWrite"

# ── Installation auto des libs optionnelles ──
def install(pkg):
    import subprocess, sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", pkg, "-q"])

try:
    import sounddevice as sd
    import soundfile as sf
    SOUND_OK = True
except ImportError:
    print("[INSTALL] Installation de sounddevice et soundfile...")
    install("sounddevice")
    install("soundfile")
    import sounddevice as sd
    import soundfile as sf
    SOUND_OK = True

try:
    from PIL import Image as PILImage
    PIL_OK = True
except ImportError:
    print("[INSTALL] Installation de Pillow...")
    install("Pillow")
    from PIL import Image as PILImage
    PIL_OK = True

try:
    import fitz  # PyMuPDF pour PDF
    PDF_OK = True
except ImportError:
    print("[INSTALL] Installation de PyMuPDF...")
    install("pymupdf")
    try:
        import fitz
        PDF_OK = True
    except ImportError:
        PDF_OK = False
        print("[WARN] PDF non disponible.")

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_OK = True
except ImportError:
    print("[INSTALL] Installation de python-pptx...")
    install("python-pptx")
    try:
        from pptx import Presentation
        PPTX_OK = True
    except ImportError:
        PPTX_OK = False
        print("[WARN] PPTX non disponible.")

# ─────────────────────────────────────────────
# TÉLÉCHARGEMENT MODÈLE MEDIAPIPE
# ─────────────────────────────────────────────

MODELS = {
    "hand_landmarker.task": (
        "https://storage.googleapis.com/mediapipe-models/"
        "hand_landmarker/hand_landmarker/float16/1/hand_landmarker.task"
    ),
}

def download_model(filename, url):
    if os.path.exists(filename):
        return True
    print(f"[DOWNLOAD] {filename}...")
    try:
        urllib.request.urlretrieve(url, filename)
        print(f"[OK] {filename} téléchargé.")
        return True
    except Exception as e:
        print(f"[ERREUR] {e}")
        return False

if not download_model("hand_landmarker.task", MODELS["hand_landmarker.task"]):
    print("[FATAL] Modèle requis. Vérifie ta connexion.")
    exit(1)

# ─────────────────────────────────────────────
# CONFIGURATION
# ─────────────────────────────────────────────

ERASER_RADIUS  = 40
PARTICLE_COUNT = 5
SCREENSHOTS_DIR = "screenshots"
RECORDINGS_DIR  = "recordings"

COLORS = {
    "Noir":   (0,     0,   0),
    "Rouge":  (0,     0, 200),
    "Bleu":   (200,  60,   0),
    "Vert":   (0,   180,  60),
    "Violet": (180,   0, 180),
    "Orange": (0,   140, 255),
    "Gomme":  None,
}

BRUSH_SIZES = [2, 5, 10, 18, 30]

# Connexions des 21 points de repère MediaPipe Hands
HAND_CONNECTIONS = [
    # Pouce
    (0, 1), (1, 2), (2, 3), (3, 4),
    # Index
    (0, 5), (5, 6), (6, 7), (7, 8),
    # Majeur
    (0, 9), (9, 10), (10, 11), (11, 12),
    # Annulaire
    (0, 13), (13, 14), (14, 15), (15, 16),
    # Auriculaire
    (0, 17), (17, 18), (18, 19), (19, 20),
    # Paume
    (5, 9), (9, 13), (13, 17),
]

# Couleurs par doigt pour les connexions (BGR)
FINGER_COLORS = {
    "thumb":  (0, 215, 255),   # or/jaune
    "index":  (0, 255, 0),     # vert
    "middle": (255, 100, 0),   # bleu clair
    "ring":   (180, 0, 255),   # violet
    "pinky":  (0, 80, 255),    # rouge-orange
    "palm":   (200, 200, 200), # gris
}

FINGER_CONNECTION_MAP = {
    (0,1):"thumb",(1,2):"thumb",(2,3):"thumb",(3,4):"thumb",
    (0,5):"index",(5,6):"index",(6,7):"index",(7,8):"index",
    (0,9):"middle",(9,10):"middle",(10,11):"middle",(11,12):"middle",
    (0,13):"ring",(13,14):"ring",(14,15):"ring",(15,16):"ring",
    (0,17):"pinky",(17,18):"pinky",(18,19):"pinky",(19,20):"pinky",
    (5,9):"palm",(9,13):"palm",(13,17):"palm",
}

for d in [SCREENSHOTS_DIR, RECORDINGS_DIR]:
    os.makedirs(d, exist_ok=True)

# ─────────────────────────────────────────────
# INIT MEDIAPIPE
# ─────────────────────────────────────────────

hand_result_global = None

def hand_callback(result, output_image, timestamp_ms):
    global hand_result_global
    hand_result_global = result

hand_options = mp_vision.HandLandmarkerOptions(
    base_options=mp_python.BaseOptions(model_asset_path="hand_landmarker.task"),
    running_mode=mp_vision.RunningMode.LIVE_STREAM,
    num_hands=1,
    min_hand_detection_confidence=0.7,
    min_hand_presence_confidence=0.7,
    min_tracking_confidence=0.5,
    result_callback=hand_callback,
)
hand_landmarker = mp_vision.HandLandmarker.create_from_options(hand_options)
print(f"[OK] {APP_NAME} — Hand Landmarker prêt.")

# ─────────────────────────────────────────────
# CLASSE PARTICULE
# ─────────────────────────────────────────────

class Particle:
    def __init__(self, x, y, color):
        self.x     = float(x)
        self.y     = float(y)
        self.vx    = np.random.uniform(-2, 2)
        self.vy    = np.random.uniform(-3, -0.5)
        self.life  = 1.0
        self.r     = int(np.random.randint(1, 4))
        self.color = color

    def update(self):
        self.x    += self.vx
        self.y    += self.vy
        self.vy   += 0.15
        self.life -= 0.08
        return self.life > 0

    def draw(self, frame):
        if self.life <= 0:
            return
        alpha = self.life
        c = tuple(max(0, min(255, int(v * alpha + 255 * (1 - alpha)))) for v in self.color)
        cv2.circle(frame, (int(self.x), int(self.y)), max(self.r, 1), c, -1, cv2.LINE_AA)


# ─────────────────────────────────────────────
# CLASSE ENREGISTREMENT AUDIO
# ─────────────────────────────────────────────

class AudioRecorder:
    SAMPLE_RATE = 44100
    CHANNELS    = 1

    def __init__(self):
        self.recording = False
        self.frames    = []
        self.thread    = None

    def start(self):
        if self.recording:
            return
        self.frames    = []
        self.recording = True
        self.thread    = threading.Thread(target=self._record, daemon=True)
        self.thread.start()
        print("[AUDIO] Enregistrement audio démarré.")

    def _record(self):
        with sd.InputStream(samplerate=self.SAMPLE_RATE,
                            channels=self.CHANNELS,
                            dtype="float32") as stream:
            while self.recording:
                data, _ = stream.read(1024)
                self.frames.append(data.copy())

    def stop(self, filepath):
        if not self.recording:
            return
        self.recording = False
        if self.thread:
            self.thread.join(timeout=2)
        if self.frames:
            audio_data = np.concatenate(self.frames, axis=0)
            sf.write(filepath, audio_data, self.SAMPLE_RATE)
            print(f"[AUDIO] Sauvegardé : {filepath}")


# ─────────────────────────────────────────────
# CLASSE CHARGEUR DE FICHIERS (PDF / PPTX)
# ─────────────────────────────────────────────

class FileLoader:
    def __init__(self):
        self.slides      = []
        self.slide_index = 0
        self.filename    = ""

    def load(self, filepath):
        ext = os.path.splitext(filepath)[1].lower()
        self.slides      = []
        self.slide_index = 0
        self.filename    = os.path.basename(filepath)

        if ext == ".pdf" and PDF_OK:
            self._load_pdf(filepath)
        elif ext in (".pptx", ".ppt") and PPTX_OK:
            self._load_pptx(filepath)
        else:
            print(f"[WARN] Format non supporté ou lib manquante : {ext}")
            return False

        print(f"[FILE] {len(self.slides)} slide(s) chargée(s) depuis {self.filename}")
        return len(self.slides) > 0

    def _load_pdf(self, filepath):
        doc = fitz.open(filepath)
        for page in doc:
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)
            img = np.frombuffer(pix.samples, dtype=np.uint8)
            img = img.reshape(pix.height, pix.width, pix.n)
            if pix.n == 4:
                img = cv2.cvtColor(img, cv2.COLOR_RGBA2BGR)
            else:
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)
            self.slides.append(img)
        doc.close()

    def _load_pptx(self, filepath):
        prs     = Presentation(filepath)
        slide_w = int(prs.slide_width.pt)
        slide_h = int(prs.slide_height.pt)

        for slide in prs.slides:
            img = np.ones((slide_h, slide_w, 3), dtype=np.uint8) * 255
            for shape in slide.shapes:
                if shape.has_text_frame:
                    x = int(shape.left.pt)
                    y = int(shape.top.pt)
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            cv2.putText(img, text, (x, y + 20),
                                        cv2.FONT_HERSHEY_SIMPLEX, 0.7,
                                        (20, 20, 20), 1, cv2.LINE_AA)
                            y += 30
            self.slides.append(img)

    def current_slide(self, w, h):
        if not self.slides:
            return None
        slide = self.slides[self.slide_index]
        return cv2.resize(slide, (w, h), interpolation=cv2.INTER_AREA)

    def next_slide(self):
        if self.slides and self.slide_index < len(self.slides) - 1:
            self.slide_index += 1

    def prev_slide(self):
        if self.slides and self.slide_index > 0:
            self.slide_index -= 1

    def has_slides(self):
        return len(self.slides) > 0

    def info(self):
        if not self.slides:
            return ""
        return f"{self.filename}  [{self.slide_index + 1}/{len(self.slides)}]"


# ─────────────────────────────────────────────
# CLASSE PRINCIPALE — AeroWrite
# ─────────────────────────────────────────────

class AeroWrite:

    def __init__(self, width=1280, height=720):
        self.w = width
        self.h = height

        self.canvas    = np.zeros((height, width, 3), dtype=np.uint8)
        self.particles = []

        self.prev_x = None
        self.prev_y = None

        self.color_name = "Noir"
        self.color      = COLORS["Noir"]
        self.brush_idx  = 1
        self.brush_size = BRUSH_SIZES[self.brush_idx]

        self.toolbar_h  = 80
        self.color_keys = list(COLORS.keys())
        self.col_w      = width // len(self.color_keys)

        # ── Mode caméra (V pour basculer) ──
        # False = tableau blanc (visage non visible)
        # True  = flux caméra en direct avec squelette main
        self.camera_mode = False

        # Enregistrement
        self.video_writer  = None
        self.audio_rec     = AudioRecorder()
        self.is_recording  = False
        self.rec_start     = 0
        self.rec_blink     = 0

        # Fichier PDF/PPTX
        self.file_loader = FileLoader()

        # Screenshot
        self.last_screenshot  = 0
        self.screenshot_flash = 0

        # FPS
        self.fps_time  = time.time()
        self.fps       = 0
        self.frame_cnt = 0

    # ─────────────────────────────────────────
    # DÉTECTION GESTES
    # ─────────────────────────────────────────

    def get_finger_state(self, landmarks):
        index_up  = landmarks[8].y  < landmarks[6].y
        middle_up = landmarks[12].y < landmarks[10].y
        ring_up   = landmarks[16].y < landmarks[14].y
        return index_up, middle_up, ring_up

    def get_index_tip(self, landmarks):
        lm = landmarks[8]
        return int(lm.x * self.w), int(lm.y * self.h)

    # ─────────────────────────────────────────
    # DESSIN SQUELETTE MAIN (21 POINTS)
    # ─────────────────────────────────────────

    def draw_hand_skeleton(self, frame, landmarks, alpha_bg=False):
        """
        Dessine les 21 points de repère + connexions MediaPipe Hands.
        alpha_bg=True → semi-transparent (mode tableau blanc)
        alpha_bg=False → opaque vif (mode caméra)
        """
        pts = []
        for lm in landmarks:
            px = int(lm.x * self.w)
            py = int(lm.y * self.h)
            pts.append((px, py))

        overlay = frame.copy() if alpha_bg else None

        # ── Connexions ──
        for (a, b), finger in FINGER_CONNECTION_MAP.items():
            color = FINGER_COLORS[finger]
            thickness = 2 if alpha_bg else 3
            target = overlay if alpha_bg else frame
            cv2.line(target, pts[a], pts[b], color, thickness, cv2.LINE_AA)

        # ── Points (cercles) ──
        for idx, (px, py) in enumerate(pts):
            target = overlay if alpha_bg else frame
            # Poignée = point 0
            if idx == 0:
                cv2.circle(target, (px, py), 10 if not alpha_bg else 7,
                           (255, 255, 255), -1, cv2.LINE_AA)
                cv2.circle(target, (px, py), 10 if not alpha_bg else 7,
                           (100, 100, 100), 2, cv2.LINE_AA)
            # Bout des doigts = 4,8,12,16,20
            elif idx in (4, 8, 12, 16, 20):
                cv2.circle(target, (px, py), 9 if not alpha_bg else 6,
                           (0, 255, 255), -1, cv2.LINE_AA)
                cv2.circle(target, (px, py), 9 if not alpha_bg else 6,
                           (0, 180, 180), 2, cv2.LINE_AA)
            else:
                cv2.circle(target, (px, py), 6 if not alpha_bg else 4,
                           (255, 255, 255), -1, cv2.LINE_AA)
                cv2.circle(target, (px, py), 6 if not alpha_bg else 4,
                           (80, 80, 80), 1, cv2.LINE_AA)

        # ── Numéros des points (mode caméra seulement) ──
        if not alpha_bg:
            for idx, (px, py) in enumerate(pts):
                cv2.putText(frame, str(idx), (px + 6, py - 4),
                            cv2.FONT_HERSHEY_SIMPLEX, 0.28, (220, 220, 220), 1, cv2.LINE_AA)

        # ── Fusion semi-transparente (mode tableau blanc) ──
        if alpha_bg and overlay is not None:
            cv2.addWeighted(overlay, 0.65, frame, 0.35, 0, frame)

    # ─────────────────────────────────────────
    # DESSIN TRAIT
    # ─────────────────────────────────────────

    def draw_stroke(self, x, y):
        if self.prev_x is None:
            self.prev_x, self.prev_y = x, y
            return
        if self.color_name == "Gomme":
            cv2.circle(self.canvas, (x, y), ERASER_RADIUS, (0, 0, 0), -1, cv2.LINE_AA)
            for _ in range(3):
                px = x + int(np.random.randint(-ERASER_RADIUS, ERASER_RADIUS))
                py = y + int(np.random.randint(-ERASER_RADIUS, ERASER_RADIUS))
                self.particles.append(Particle(px, py, (200, 200, 200)))
        else:
            cv2.line(self.canvas,
                     (self.prev_x, self.prev_y), (x, y),
                     self.color, self.brush_size, cv2.LINE_AA)
            for _ in range(PARTICLE_COUNT):
                px = x + int(np.random.randint(-self.brush_size, self.brush_size + 1))
                py = y + int(np.random.randint(-self.brush_size, self.brush_size + 1))
                self.particles.append(Particle(px, py, self.color))
        self.prev_x, self.prev_y = x, y

    def update_particles(self, frame):
        self.particles = [p for p in self.particles if p.update()]
        for p in self.particles:
            p.draw(frame)

    def clear_canvas(self):
        self.canvas[:] = 0
        self.particles.clear()
        self.prev_x = self.prev_y = None

    # ─────────────────────────────────────────
    # ENREGISTREMENT VIDÉO + AUDIO
    # ─────────────────────────────────────────

    def start_recording(self, w, h):
        if self.is_recording:
            return
        ts              = time.strftime("%Y%m%d_%H%M%S")
        self.video_path = os.path.join(RECORDINGS_DIR, f"aerowrite_{ts}.mp4")
        self.audio_path = os.path.join(RECORDINGS_DIR, f"aerowrite_{ts}.wav")
        fourcc = cv2.VideoWriter_fourcc(*"mp4v")
        self.video_writer = cv2.VideoWriter(self.video_path, fourcc, 20, (w, h))
        self.audio_rec.start()
        self.is_recording = True
        self.rec_start    = time.time()
        print(f"[REC] Début : {self.video_path}")

    def stop_recording(self):
        if not self.is_recording:
            return
        self.is_recording = False
        if self.video_writer:
            self.video_writer.release()
            self.video_writer = None
        self.audio_rec.stop(self.audio_path)
        print(f"[REC] Fin. Vidéo : {self.video_path}")
        print(f"[REC] Audio  : {self.audio_path}")
        ts = os.path.splitext(os.path.basename(self.video_path))[0]
        print(f"[REC] Fusionner : ffmpeg -i {self.video_path} -i {self.audio_path} "
              f"-c:v copy -c:a aac {RECORDINGS_DIR}/{ts}_final.mp4")

    def write_frame(self, frame):
        if self.is_recording and self.video_writer:
            self.video_writer.write(frame)

    # ─────────────────────────────────────────
    # TOOLBAR
    # ─────────────────────────────────────────

    def draw_toolbar(self, frame):
        overlay = frame.copy()
        cv2.rectangle(overlay, (0, 0), (self.w, self.toolbar_h), (235, 235, 235), -1)
        cv2.addWeighted(overlay, 0.92, frame, 0.08, 0, frame)
        cv2.line(frame, (0, self.toolbar_h), (self.w, self.toolbar_h),
                 (180, 180, 180), 1)

        # Cases couleur
        for i, name in enumerate(self.color_keys):
            x1 = i * self.col_w
            x2 = x1 + self.col_w

            is_active = (name == self.color_name)
            if name == "Gomme":
                color_bgr = (220, 220, 220)
                border    = (120, 120, 120)
            else:
                color_bgr = COLORS[name]
                border    = (80, 80, 80)

            cv2.rectangle(frame, (x1 + 6, 10), (x2 - 6, self.toolbar_h - 10),
                          color_bgr, -1, cv2.LINE_AA)
            cv2.rectangle(frame, (x1 + 6, 10), (x2 - 6, self.toolbar_h - 10),
                          border, 1, cv2.LINE_AA)
            if is_active:
                cv2.rectangle(frame, (x1 + 3, 7), (x2 - 3, self.toolbar_h - 7),
                              (60, 60, 200), 2, cv2.LINE_AA)

            lbl_color = (255, 255, 255) if name not in ("Gomme", "Jaune", "Orange") else (40, 40, 40)
            cv2.putText(frame, name[:5], (x1 + 10, self.toolbar_h // 2 + 5),
                        cv2.FONT_HERSHEY_SIMPLEX, 0.35, lbl_color, 1, cv2.LINE_AA)

        # Taille pinceau
        cv2.putText(frame, f"Taille: {self.brush_size}px",
                    (self.w - 420, 28),
                    cv2.FONT_HERSHEY_SIMPLEX, 0.45, (60, 60, 60), 1, cv2.LINE_AA)

        # Mode actuel
        mode_txt   = "CAMERA" if self.camera_mode else "TABLEAU BLANC"
        mode_color = (30, 120, 200) if self.camera_mode else (40, 140, 40)
        cv2.putText(frame, f"[V] Mode: {mode_txt}",
                    (self.w - 420, 52),
                    cv2.FONT_HERSHEY_SIMPLEX, 0.42, mode_color, 1, cv2.LINE_AA)

        # Raccourcis clavier
        cv2.putText(frame,
                    "[R]ec  [O]uvrir  [<][>]slides  [C]lear  [S]creen  [+][-]  [V]cam  [Q]uit",
                    (self.w - 580, self.toolbar_h - 10),
                    cv2.FONT_HERSHEY_SIMPLEX, 0.33, (100, 100, 100), 1, cv2.LINE_AA)

    def check_toolbar_click(self, x, y):
        if y > self.toolbar_h:
            return
        idx = x // self.col_w
        if 0 <= idx < len(self.color_keys):
            self.color_name = self.color_keys[idx]
            self.color      = COLORS[self.color_name]
            self.prev_x     = self.prev_y = None

    # ─────────────────────────────────────────
    # HUD
    # ─────────────────────────────────────────

    def draw_hud(self, frame, drawing_mode):
        self.frame_cnt += 1
        now = time.time()
        if now - self.fps_time >= 1.0:
            self.fps       = self.frame_cnt
            self.frame_cnt = 0
            self.fps_time  = now

        y_base = self.toolbar_h + 28

        mode_label = "DESSIN" if drawing_mode else "NAVIGATION"
        mode_color = (0, 150, 0) if drawing_mode else (150, 100, 0)
        cv2.putText(frame, f"Mode: {mode_label}", (10, y_base),
                    cv2.FONT_HERSHEY_SIMPLEX, 0.55, mode_color, 2, cv2.LINE_AA)

        cv2.putText(frame, f"FPS: {self.fps}", (10, y_base + 26),
                    cv2.FONT_HERSHEY_SIMPLEX, 0.45, (120, 120, 120), 1, cv2.LINE_AA)

        if self.file_loader.has_slides():
            info = self.file_loader.info()
            cv2.putText(frame, info,
                        (self.w // 2 - len(info) * 4, y_base),
                        cv2.FONT_HERSHEY_SIMPLEX, 0.5, (40, 40, 160), 1, cv2.LINE_AA)

        if self.is_recording:
            elapsed   = int(now - self.rec_start)
            m, s      = divmod(elapsed, 60)
            self.rec_blink = (self.rec_blink + 1) % 30
            if self.rec_blink < 15:
                cv2.circle(frame, (self.w - 20, y_base - 6), 8, (0, 0, 220), -1, cv2.LINE_AA)
            cv2.putText(frame, f"REC {m:02d}:{s:02d}",
                        (self.w - 110, y_base),
                        cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 0, 200), 2, cv2.LINE_AA)

        if self.screenshot_flash > 0:
            wh = np.ones_like(frame) * 255
            a  = self.screenshot_flash / 8
            cv2.addWeighted(wh, a * 0.5, frame, 1 - a * 0.5, 0, frame)
            cv2.putText(frame, "CAPTURE !",
                        (self.w // 2 - 80, self.h // 2),
                        cv2.FONT_HERSHEY_SIMPLEX, 1.4, (0, 150, 0), 3, cv2.LINE_AA)
            self.screenshot_flash -= 1

        cv2.putText(frame, "☝ Index = dessin  |  ✌ Index+Majeur = navigation",
                    (10, self.h - 12),
                    cv2.FONT_HERSHEY_SIMPLEX, 0.4, (150, 150, 150), 1, cv2.LINE_AA)

    # ─────────────────────────────────────────
    # CURSEUR INDEX
    # ─────────────────────────────────────────

    def draw_cursor(self, frame, x, y):
        if self.color_name == "Gomme":
            cv2.circle(frame, (x, y), ERASER_RADIUS, (150, 150, 150), 2, cv2.LINE_AA)
            cv2.line(frame, (x - ERASER_RADIUS, y), (x + ERASER_RADIUS, y),
                     (180, 180, 180), 1, cv2.LINE_AA)
            cv2.line(frame, (x, y - ERASER_RADIUS), (x, y + ERASER_RADIUS),
                     (180, 180, 180), 1, cv2.LINE_AA)
        else:
            cv2.circle(frame, (x, y), self.brush_size + 4, self.color, 2, cv2.LINE_AA)
            cv2.circle(frame, (x, y), 3, (0, 0, 0), -1, cv2.LINE_AA)

    # ─────────────────────────────────────────
    # SCREENSHOT
    # ─────────────────────────────────────────

    def take_screenshot(self, frame):
        now = time.time()
        if now - self.last_screenshot < 2.0:
            return
        self.last_screenshot  = now
        self.screenshot_flash = 8
        fname = os.path.join(SCREENSHOTS_DIR,
                             f"aerowrite_{time.strftime('%Y%m%d_%H%M%S')}.png")
        cv2.imwrite(fname, frame)
        print(f"[CAPTURE] {fname}")

    # ─────────────────────────────────────────
    # OUVRIR FICHIER
    # ─────────────────────────────────────────

    def open_file_dialog(self):
        try:
            import tkinter as tk
            from tkinter import filedialog
            root = tk.Tk()
            root.withdraw()
            root.attributes("-topmost", True)
            path = filedialog.askopenfilename(
                title="AeroWrite — Ouvrir un fichier",
                filetypes=[
                    ("Fichiers supportés", "*.pdf *.pptx *.ppt"),
                    ("PDF", "*.pdf"),
                    ("PowerPoint", "*.pptx *.ppt"),
                ],
            )
            root.destroy()
            return path
        except Exception as e:
            print(f"[WARN] Dialogue impossible : {e}")
            return None

    # ─────────────────────────────────────────
    # BOUCLE PRINCIPALE
    # ─────────────────────────────────────────

    def run(self):
        global hand_result_global

        cap = cv2.VideoCapture(0)
        cap.set(cv2.CAP_PROP_FRAME_WIDTH,  self.w)
        cap.set(cv2.CAP_PROP_FRAME_HEIGHT, self.h)

        if not cap.isOpened():
            print("[ERREUR] Webcam introuvable.")
            return

        actual_w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        actual_h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        if actual_w != self.w or actual_h != self.h:
            self.w      = actual_w
            self.h      = actual_h
            self.canvas = np.zeros((self.h, self.w, 3), dtype=np.uint8)
            self.col_w  = self.w // len(self.color_keys)

        print("=" * 65)
        print(f"   ✈  {APP_NAME} — Tableau Blanc Interactif  ✈")
        print("=" * 65)
        print("  Gestes :")
        print("    ☝  Index seul       → DESSIN")
        print("    ✌  Index + Majeur   → NAVIGATION / toolbar")
        print("  Clavier :")
        print("    V   → Basculer Mode Caméra / Tableau Blanc")
        print("    R   → Démarrer / Arrêter enregistrement")
        print("    O   → Ouvrir PDF ou PowerPoint")
        print("    <   → Slide précédente")
        print("    >   → Slide suivante")
        print("    C   → Effacer le tableau")
        print("    S   → Screenshot")
        print("    +/- → Taille pinceau")
        print("    Q   → Quitter")
        print("=" * 65)
        print(f"  Démarrage en mode : TABLEAU BLANC (appuie sur V pour la caméra)")
        print("=" * 65)

        timestamp = 0

        while True:
            ret, frame = cap.read()
            if not ret:
                break

            frame = cv2.flip(frame, 1)
            h_f, w_f = frame.shape[:2]
            if w_f != self.w or h_f != self.h:
                self.w      = w_f
                self.h      = h_f
                self.canvas = np.zeros((self.h, self.w, 3), dtype=np.uint8)
                self.col_w  = self.w // len(self.color_keys)

            # ── MediaPipe ────────────────────────────
            rgb      = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            mp_image = mp.Image(image_format=mp.ImageFormat.SRGB, data=rgb)
            timestamp += 1
            hand_landmarker.detect_async(mp_image, timestamp)

            drawing_mode       = False
            ix, iy             = None, None
            landmarks_detected = None

            if hand_result_global and hand_result_global.hand_landmarks:
                landmarks          = hand_result_global.hand_landmarks[0]
                landmarks_detected = landmarks
                ix, iy             = self.get_index_tip(landmarks)
                index_up, middle_up, ring_up = self.get_finger_state(landmarks)

                if index_up and not middle_up:
                    drawing_mode = True
                    if iy > self.toolbar_h:
                        self.draw_stroke(ix, iy)
                    else:
                        self.check_toolbar_click(ix, iy)
                        self.prev_x = self.prev_y = None
                elif index_up and middle_up:
                    self.prev_x = self.prev_y = None
                    if iy <= self.toolbar_h:
                        self.check_toolbar_click(ix, iy)
                else:
                    self.prev_x = self.prev_y = None

            # ─────────────────────────────────────────
            # CONSTRUCTION DU FOND SELON LE MODE
            # ─────────────────────────────────────────

            if self.camera_mode:
                # ── MODE CAMÉRA : fond = flux webcam en direct ──
                display = frame.copy()

                # Superposer dessin
                canvas_gray = cv2.cvtColor(self.canvas, cv2.COLOR_BGR2GRAY)
                _, mask     = cv2.threshold(canvas_gray, 1, 255, cv2.THRESH_BINARY)
                draw_layer  = cv2.bitwise_and(self.canvas, self.canvas, mask=mask)
                display[mask > 0] = draw_layer[mask > 0]

                # Squelette main 21 points (vif, opaque)
                if landmarks_detected:
                    self.draw_hand_skeleton(display, landmarks_detected, alpha_bg=False)

                # Badge mode caméra
                cv2.rectangle(display, (self.w - 160, self.toolbar_h + 4),
                              (self.w - 4, self.toolbar_h + 26), (0, 80, 180), -1)
                cv2.putText(display, "MODE CAMERA",
                            (self.w - 155, self.toolbar_h + 20),
                            cv2.FONT_HERSHEY_SIMPLEX, 0.42, (255, 255, 255), 1, cv2.LINE_AA)

            else:
                # ── MODE TABLEAU BLANC : fond blanc ou slide (visage NON visible) ──
                if self.file_loader.has_slides():
                    display = self.file_loader.current_slide(self.w, self.h).copy()
                else:
                    display = np.ones((self.h, self.w, 3), dtype=np.uint8) * 255

                # Superposer dessin
                canvas_gray = cv2.cvtColor(self.canvas, cv2.COLOR_BGR2GRAY)
                _, mask     = cv2.threshold(canvas_gray, 1, 255, cv2.THRESH_BINARY)
                draw_layer  = cv2.bitwise_and(self.canvas, self.canvas, mask=mask)
                display[mask > 0] = draw_layer[mask > 0]

                # Squelette main 21 points (semi-transparent pour ne pas gêner)
                if landmarks_detected:
                    self.draw_hand_skeleton(display, landmarks_detected, alpha_bg=True)

            # ── Éléments communs ──────────────────────
            self.update_particles(display)
            self.draw_toolbar(display)
            if ix is not None:
                self.draw_cursor(display, ix, iy)
            self.draw_hud(display, drawing_mode)

            # Enregistrement
            self.write_frame(display)

            cv2.imshow(APP_NAME, display)

            # ── Clavier ───────────────────────────────
            key = cv2.waitKey(1) & 0xFF

            if key in (ord('q'), 27):
                break

            elif key == ord('v'):
                self.camera_mode = not self.camera_mode
                mode_str = "CAMÉRA" if self.camera_mode else "TABLEAU BLANC"
                print(f"[MODE] Basculé en → {mode_str}")

            elif key == ord('r'):
                if not self.is_recording:
                    self.start_recording(self.w, self.h)
                else:
                    self.stop_recording()

            elif key == ord('o'):
                path = self.open_file_dialog()
                if path:
                    self.clear_canvas()
                    ok = self.file_loader.load(path)
                    if not ok:
                        print("[WARN] Impossible de charger ce fichier.")

            elif key in (ord('<'), ord(','), 81, 2):
                self.file_loader.prev_slide()
                self.clear_canvas()

            elif key in (ord('>'), ord('.'), 83, 3):
                self.file_loader.next_slide()
                self.clear_canvas()

            elif key == ord('c'):
                self.clear_canvas()

            elif key == ord('s'):
                self.take_screenshot(display)

            elif key in (ord('+'), ord('=')):
                self.brush_idx  = min(self.brush_idx + 1, len(BRUSH_SIZES) - 1)
                self.brush_size = BRUSH_SIZES[self.brush_idx]

            elif key == ord('-'):
                self.brush_idx  = max(self.brush_idx - 1, 0)
                self.brush_size = BRUSH_SIZES[self.brush_idx]

        # ── Nettoyage ──────────────────────────────
        if self.is_recording:
            self.stop_recording()
        cap.release()
        cv2.destroyAllWindows()
        hand_landmarker.close()
        print(f"[{APP_NAME}] Fermé. À bientôt !")


# ─────────────────────────────────────────────
if __name__ == "__main__":
    app = AeroWrite(width=1280, height=720)
    app.run()